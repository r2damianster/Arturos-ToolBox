"""
logic/utilidades.py
Lógica de todas las utilidades adaptada para contexto web (Flask).
"""

import io
import os
import csv
import re
import zipfile
import random
import math

import qrcode
from PIL import Image
from fpdf import FPDF
import pandas as pd
import numpy as np
from scipy import stats
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────
# 1. CSV → BIB
# ─────────────────────────────────────────────

def csv_a_bib(csv_bytes: bytes) -> bytes:
    contenido = csv_bytes.decode("utf-8")
    reader = csv.DictReader(io.StringIO(contenido))
    registros = list(reader)

    campos = {
        "title": "title",
        "author": "authors",
        "year": "year",
        "journal": "journal",
        "volume": "volume",
        "number": "issue",
        "pages": "pages",
        "doi": "doi",
        "url": "url",
        "note": "notes",
    }

    salida = io.StringIO()
    for i, row in enumerate(registros):
        bib_id = row.get("key", f"art_{i}")
        salida.write(f"@article{{{bib_id},\n")
        for bib_field, csv_field in campos.items():
            valor = row.get(csv_field, "").strip()
            if valor:
                valor = valor.replace("&", "\\&").replace("_", "\\_")
                salida.write(f"  {bib_field} = {{{valor}}},\n")
        salida.write("}\n\n")

    return salida.getvalue().encode("utf-8")


# ─────────────────────────────────────────────
# 2. QR
# ─────────────────────────────────────────────

def generar_qr(texto: str) -> bytes:
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(texto)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    # Convert to RGB if necessary (for JPEG format)
    if img.mode not in ('RGB', 'L'):
        img = img.convert('RGB')
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=95)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# 3. Slides (Markdown → PPTX)
# ─────────────────────────────────────────────

def _aplicar_formato_markdown(parrafo, texto_linea: str, tamaño_base: int, color_rgb: tuple):
    url_pattern = r"(https?://[^\s]+)"
    partes_negrita = re.split(r"(\*.*?\*)", texto_linea)

    for parte in partes_negrita:
        if not parte:
            continue
        sub_partes = re.split(url_pattern, parte)
        es_negrita = parte.startswith("*") and parte.endswith("*")

        for sub in sub_partes:
            if not sub:
                continue
            run = parrafo.add_run()
            run.font.name = "Segoe UI"
            run.font.size = Pt(tamaño_base)

            if re.match(url_pattern, sub):
                run.text = sub
                run.font.color.rgb = RGBColor(0, 102, 204)
                run.font.underline = True
                run.hyperlink.address = sub
            elif es_negrita:
                run.text = sub.replace("*", "")
                run.font.bold = True
                run.font.color.rgb = RGBColor(*color_rgb)
            else:
                run.text = re.sub(r"^[-*]\s*", "", sub)
                run.font.color.rgb = RGBColor(*color_rgb)


def _añadir_footer(slide, color_rgb: tuple):
    rect = slide.shapes.add_shape(1, 0, Inches(7.1), Inches(10), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*color_rgb)
    rect.line.visible = False


def generar_pptx(lineas: list, template_path: str = None) -> bytes:
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    AZUL   = (44, 62, 80)
    CELESTE = (174, 214, 241)
    GRIS   = (60, 60, 60)
    slide_actual = None

    for linea in lineas:
        linea_limpia = linea.strip()
        if not linea_limpia:
            continue

        if linea_limpia.startswith("##"):
            texto_titulo = linea_limpia.replace("##", "").strip()
            slide_actual = prs.slides.add_slide(prs.slide_layouts[6])
            _añadir_footer(slide_actual, CELESTE)

            tb = slide_actual.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = texto_titulo
            p.font.bold = True
            p.font.size = Pt(32)
            p.font.color.rgb = RGBColor(*AZUL)
            p.alignment = PP_ALIGN.LEFT

        elif slide_actual:
            if len(slide_actual.shapes) < 3:
                body = slide_actual.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
                body.text_frame.word_wrap = True
                tf = body.text_frame
            else:
                tf = slide_actual.shapes[-1].text_frame

            espacios = len(linea) - len(linea.lstrip())
            nivel = 1 if espacios >= 4 else 0
            p = tf.add_paragraph()
            p.level = nivel
            p.space_before = Pt(12)
            _aplicar_formato_markdown(p, linea_limpia, 20 if nivel == 0 else 16, GRIS)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# 4. Unir imágenes → PDF
# ─────────────────────────────────────────────

def unir_imagenes_a_pdf(archivos_imagen: list) -> bytes:
    imagenes = []
    for f in archivos_imagen:
        img = Image.open(f.stream)
        if img.mode != "RGB":
            img = img.convert("RGB")
        imagenes.append(img)

    if not imagenes:
        raise ValueError("No se proporcionaron imágenes.")

    buf = io.BytesIO()
    imagenes[0].save(buf, format="PDF", save_all=True, append_images=imagenes[1:], resolution=100.0)
    buf.seek(0)
    return buf.read()


def unir_imagenes_a_jpg(archivos_imagen: list, orientacion: str = "vertical") -> bytes:
    """Une varias imágenes en una sola imagen JPG."""
    imagenes = []
    for f in archivos_imagen:
        img = Image.open(f.stream)
        if img.mode != "RGB":
            img = img.convert("RGB")
        imagenes.append(img)

    if not imagenes:
        raise ValueError("No se proporcionaron imágenes.")

    if len(imagenes) == 1:
        buf = io.BytesIO()
        imagenes[0].save(buf, format="JPEG", quality=95)
        buf.seek(0)
        return buf.read()

    if orientacion == "horizontal":
        ancho_total = sum(img.width for img in imagenes)
        alto_maximo = max(img.height for img in imagenes)
        resultado = Image.new("RGB", (ancho_total, alto_maximo), color="white")
        x_offset = 0
        for img in imagenes:
            y_offset = (alto_maximo - img.height) // 2
            resultado.paste(img, (x_offset, y_offset))
            x_offset += img.width
    else:  # vertical por defecto
        ancho_maximo = max(img.width for img in imagenes)
        alto_total = sum(img.height for img in imagenes)
        resultado = Image.new("RGB", (ancho_maximo, alto_total), color="white")
        y_offset = 0
        for img in imagenes:
            x_offset = (ancho_maximo - img.width) // 2
            resultado.paste(img, (x_offset, y_offset))
            y_offset += img.height

    buf = io.BytesIO()
    resultado.save(buf, format="JPEG", quality=95)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# 5. Aplanar archivos → ZIP
# ─────────────────────────────────────────────

def aplanar_archivos(archivos: list) -> bytes:
    buf = io.BytesIO()
    nombres_vistos = {}

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in archivos:
            nombre_base = os.path.basename(f.filename) or f.filename
            if nombre_base in nombres_vistos:
                nombres_vistos[nombre_base] += 1
                nombre, ext = os.path.splitext(nombre_base)
                nombre_base = f"{nombre}_{nombres_vistos[nombre_base]}{ext}"
            else:
                nombres_vistos[nombre_base] = 0
            zf.writestr(f"archivos_extraidos/{nombre_base}", f.read())

    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# 6. Generar archivos Excel → ZIP
# ─────────────────────────────────────────────

def generar_excels(tipo: str, cantidad: int) -> bytes:
    buf_zip = io.BytesIO()

    with zipfile.ZipFile(buf_zip, "w", zipfile.ZIP_DEFLATED) as zf:
        for i in range(1, cantidad + 1):
            if tipo == "1":
                data = {
                    "Genero": [random.choice([1, 2]) for _ in range(385)],
                    "Edad": [random.randint(25, 55) for _ in range(385)],
                    "Habilidades de enseñanza": [random.randint(1, 5) for _ in range(385)],
                    "Métodos de evaluación": [random.randint(1, 5) for _ in range(385)],
                    "Uso de recursos educativos": [random.randint(1, 5) for _ in range(385)],
                    "Relaciones interpersonales": [random.randint(1, 5) for _ in range(385)],
                    "Gestión del aula": [random.randint(1, 5) for _ in range(385)],
                    "Resolución de conflictos": [random.randint(1, 5) for _ in range(385)],
                    "Adaptaciones curriculares": [random.randint(1, 5) for _ in range(385)],
                    "Atención a la diversidad cultural": [random.randint(1, 5) for _ in range(385)],
                    "Apoyo a estudiantes con necesidades especiales": [random.randint(1, 5) for _ in range(385)],
                }
                nombre = f"articulo_cientifico_{i}.xlsx"
            else:
                ejes = {
                    "Habilidades de comunicación": ["Escucha activa", "Empatía", "Expresión oral", "Turnos de palabra", "Claridad al hablar"],
                    "Lectura": ["Comprensión lectora", "Velocidad de lectura", "Interpretación de textos", "Vocabulario", "Identificación de ideas principales"],
                    "Motricidad": ["Coordinación mano-ojo", "Precisión motriz", "Fuerza en agarre", "Control de movimientos finos", "Equilibrio"],
                    "Lenguaje simbólico": ["Uso de símbolos", "Reconocimiento de patrones", "Interpretación de gráficos", "Lenguaje matemático", "Uso de mapas conceptuales"],
                    "Escritura": ["Caligrafía", "Ortografía", "Redacción", "Uso de signos de puntuación", "Coherencia en textos"],
                }
                data = {
                    "Dimensión": [eje for eje in ejes for _ in range(5)],
                    "Indicador": [ind for inds in ejes.values() for ind in inds],
                    "PreTest": [random.randint(1, 5) for _ in range(25)],
                    "PostTest": [random.randint(1, 5) for _ in range(25)],
                }
                nombre = f"caso_estudio_{i}.xlsx"

            df = pd.DataFrame(data)
            excel_buf = io.BytesIO()
            df.to_excel(excel_buf, index=False)
            zf.writestr(nombre, excel_buf.getvalue())

    buf_zip.seek(0)
    return buf_zip.read()


# ─────────────────────────────────────────────
# 7. Crear estructura de carpetas → ZIP
# ─────────────────────────────────────────────

def crear_estructura_carpetas(nombres: list) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        carpetas_creadas = set()
        
        for linea in nombres:
            linea = linea.strip()
            if not linea:
                continue
            
            # Soporta sintaxis: "carpeta/subcarpeta/subsubcarpeta"
            partes = [p.strip() for p in linea.replace("\\", "/").split("/") if p.strip()]
            
            ruta_acumulada = ""
            for i, parte in enumerate(partes):
                ruta_anterior = ruta_acumulada
                ruta_acumulada = f"{ruta_acumulada}/{parte}" if ruta_acumulada else parte
                
                # Evita crear duplicados
                if ruta_acumulada not in carpetas_creadas:
                    zf.writestr(f"{ruta_acumulada}/.gitkeep", "")
                    carpetas_creadas.add(ruta_acumulada)
    
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# 8. Convertir archivos de texto → ZIP de PDFs
# ─────────────────────────────────────────────

def convertir_a_pdf(archivos: list) -> bytes:
    buf_zip = io.BytesIO()

    with zipfile.ZipFile(buf_zip, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in archivos:
            contenido = f.read().decode("utf-8", errors="replace")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for linea in contenido.splitlines():
                pdf.cell(0, 8, txt=linea[:200], ln=True)

            nombre_pdf = f"PDF_{os.path.splitext(f.filename)[0]}.pdf"
            pdf_output = pdf.output(dest="S").encode("latin-1", errors="ignore")
            zf.writestr(nombre_pdf, pdf_output)

    buf_zip.seek(0)
    return buf_zip.read()


# ─────────────────────────────────────────────
# 9. Reducir imágenes → ZIP
# ─────────────────────────────────────────────

def reducir_imagenes(archivos: list, max_px: int = 1280, calidad: int = 60) -> bytes:
    buf_zip = io.BytesIO()

    with zipfile.ZipFile(buf_zip, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in archivos:
            img = Image.open(f.stream)
            img = img.convert("RGB")
            img.thumbnail((max_px, max_px), Image.LANCZOS)

            img_buf = io.BytesIO()
            img.save(img_buf, format="JPEG", quality=calidad)
            nombre = f"reducida_{os.path.splitext(f.filename)[0]}.jpg"
            zf.writestr(nombre, img_buf.getvalue())

    buf_zip.seek(0)
    return buf_zip.read()


# ─────────────────────────────────────────────
# 10. Análisis Pre-test / Post-test
# ─────────────────────────────────────────────

def analizar_pretest_posttest(pretest: list, posttest: list) -> dict:
    pre  = np.array(pretest, dtype=float)
    post = np.array(posttest, dtype=float)
    diff = post - pre
    n    = len(diff)

    # ── Shapiro-Wilk sobre las diferencias ──────────────────────────────
    if n >= 3:
        sw_stat, sw_p = stats.shapiro(diff)
    else:
        sw_stat, sw_p = 1.0, 1.0
    diferencias_normales = bool(sw_p >= 0.05)

    # ── Wilcoxon Signed-Rank ─────────────────────────────────────────────
    stat_w, p_wilcoxon = stats.wilcoxon(pre, post, alternative="two-sided")

    nonzero_diff = diff[diff != 0]
    ranks_abs    = stats.rankdata(np.abs(nonzero_diff))
    pos_ranks    = int(np.sum(ranks_abs[nonzero_diff > 0]))
    neg_ranks    = int(np.sum(ranks_abs[nonzero_diff < 0]))
    ties         = int(np.sum(diff == 0))
    n_eff        = n - ties

    z_approx = float((stat_w - (n_eff * (n_eff + 1)) / 4) /
                     math.sqrt((n_eff * (n_eff + 1) * (2 * n_eff + 1)) / 24))
    r_effect = abs(z_approx) / math.sqrt(n_eff) if n_eff > 0 else 0.0

    # Interpretación tamaño del efecto (r)
    if r_effect < 0.1:
        efecto_label = "insignificante"
    elif r_effect < 0.3:
        efecto_label = "pequeño"
    elif r_effect < 0.5:
        efecto_label = "mediano"
    else:
        efecto_label = "grande"

    # ── T-Student pareada ────────────────────────────────────────────────
    t_stat, p_t = stats.ttest_rel(pre, post)
    media_diff  = float(np.mean(diff))
    sd_diff     = float(np.std(diff, ddof=1))
    se          = sd_diff / math.sqrt(n)
    t_critico   = float(stats.t.ppf(0.975, df=n - 1))
    margen      = t_critico * se
    ci_lower    = media_diff - margen
    ci_upper    = media_diff + margen
    cohen_d     = media_diff / sd_diff if sd_diff > 0 else 0.0

    # Interpretación d de Cohen
    d_abs = abs(cohen_d)
    if d_abs < 0.2:
        cohen_label = "insignificante"
    elif d_abs < 0.5:
        cohen_label = "pequeño"
    elif d_abs < 0.8:
        cohen_label = "mediano"
    else:
        cohen_label = "grande"

    # ── Recomendación de test según normalidad y N ───────────────────────
    if n < 10:
        recomendacion = "muestras_pequeñas"
    elif diferencias_normales:
        recomendacion = "t_student"
    else:
        recomendacion = "wilcoxon"

    return {
        "n": n,
        "media_pre":  round(float(np.mean(pre)),  4),
        "media_post": round(float(np.mean(post)), 4),
        "datos_pre":  pre.tolist(),
        "datos_post": post.tolist(),
        "diferencias": diff.tolist(),
        "normalidad": {
            "shapiro_W":  round(float(sw_stat), 4),
            "shapiro_p":  round(float(sw_p), 4),
            "normal":     diferencias_normales,
            "recomendacion": recomendacion,
        },
        "wilcoxon": {
            "W":             round(float(stat_w), 4),
            "z":             round(z_approx, 4),
            "p_value":       round(float(p_wilcoxon), 4),
            "effect_size_r": round(r_effect, 4),
            "efecto_label":  efecto_label,
            "n_efectivo":    n_eff,
            "rangos_pos":    pos_ranks,
            "rangos_neg":    neg_ranks,
            "empates":       ties,
            "significativo": bool(p_wilcoxon < 0.05),
        },
        "t_student": {
            "t_value":       round(float(t_stat), 4),
            "p_value":       round(float(p_t), 4),
            "gl":            n - 1,
            "cohen_d":       round(cohen_d, 4),
            "cohen_label":   cohen_label,
            "significativo": bool(p_t < 0.05),
        },
        "intervalo_confianza": {
            "media_diferencia": round(media_diff, 4),
            "desviacion_std":   round(sd_diff, 4),
            "error_estandar":   round(se, 4),
            "t_critico":        round(t_critico, 4),
            "margen_error":     round(margen, 4),
            "ic_inferior":      round(ci_lower, 4),
            "ic_superior":      round(ci_upper, 4),
            "nivel_confianza":  "95%",
        },
    }
